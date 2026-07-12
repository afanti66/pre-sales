#!/usr/bin/env python3
"""生成项目PPT演示文稿"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import json, os
from datetime import datetime

OUTPUT = "/home/afanti/projects/sales-agent-dashboard/nfs-data/projects/user-management/artifacts/ppt/project-presentation.pptx"
NFS = "/home/afanti/projects/sales-agent-dashboard/nfs-data/projects/user-management"
now = datetime.now().strftime('%Y年%m月%d日')

prs = Presentation()
prs.slide_width = Inches(13.333)  # 16:9
prs.slide_height = Inches(7.5)

# 颜色方案
BLUE = RGBColor(0x1A, 0x73, 0xE8)
DARK = RGBColor(0x20, 0x20, 0x20)
GRAY = RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFA)

def add_bg(slide, color=BLUE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False, color=DARK, alignment=PP_ALIGN.LEFT, font_name='Microsoft YaHei'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_slide(slide, items, left=0.8, top=1.8, width=11.5, font_size=16, color=DARK):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Microsoft YaHei'
        p.space_after = Pt(8)
        p.level = 0

# ====== Slide 1: 封面 ======
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide)
add_text_box(slide, 1, 2, 11, 1.2, '统一用户管理与身份认证系统', 40, True, WHITE, PP_ALIGN.CENTER)
add_text_box(slide, 1, 3.3, 11, 0.8, '技术方案汇报', 28, False, WHITE, PP_ALIGN.CENTER)
add_text_box(slide, 1, 5.0, 11, 0.5, f'某单位  |  {now}', 18, False, RGBColor(0xBB,0xDE,0xFB), PP_ALIGN.CENTER)

# ====== Slide 2: 现状与痛点 ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_text_box(slide, 0.8, 0.5, 11, 0.8, '项目背景与痛点', 32, True, BLUE)
add_bullet_slide(slide, [
    '▸ 多系统独立管理账号：OA、财务、人事、档案、项目管理——5套账号，5套密码',
    '▸ 人员变动时无法统一禁用：离职人员账号仍可访问部分系统，存在安全隐患',
    '▸ 权限模型不一致：各系统权限体系各自为政，审计无法统一',
    '▸ 审计日志分散：无法实现全局安全审计，合规检查耗时耗力',
    '',
    '核心矛盾：', 
    '   50,000+ 用户  ×  8+ 业务系统  =  安全管控盲区'
])

# ====== Slide 3: 建设目标 ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_text_box(slide, 0.8, 0.5, 11, 0.8, '建设目标', 32, True, BLUE)
add_bullet_slide(slide, [
    '目标一：统一用户管理 —— 组织架构、用户账号、角色权限的集中管理',
    '目标二：统一身份认证 —— SSO单点登录，支持多因子+国密',
    '目标三：统一权限管理 —— RBAC+ABAC细粒度权限控制',
    '目标四：统一安全审计 —— 全平台操作日志集中采集、存储、分析',
    '目标五：统一应用接入 —— 标准化接口，新系统快速接入',
    '',
    '系统规模：50,000用户  |  500 TPS并发  |  99.9%可用性'
])

# ====== Slide 4: 总体架构 ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_text_box(slide, 0.8, 0.5, 11, 0.8, '总体架构', 32, True, BLUE)
add_bullet_slide(slide, [
    '展示层：Vue 3 + Element Plus 管理控制台',
    '接入层：Spring Cloud Gateway（路由/限流/鉴权/日志）',
    '业务层（微服务）：',
    '   用户管理服务(8081) → 组织/用户/同步',
    '   身份认证服务(8082) → 认证/SSO/令牌',
    '   权限管理服务(8083) → 角色/授权/数据权限',
    '   安全审计服务(8084) → 日志/检索/告警',
    '数据层：MySQL MGR + Redis Cluster + ES Cluster',
    '通信：同步(OpenFeign) + 异步(RabbitMQ)'
])

# ====== Slide 5: 技术方案亮点 ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_text_box(slide, 0.8, 0.5, 11, 0.8, '技术方案亮点', 32, True, BLUE)
add_bullet_slide(slide, [
    '① 多因子认证体系：密码 + 短信 + TOTP + SM2证书，分级安全策略',
    '② 全协议SSO兼容：OAuth2 / OIDC / SAML2.0 / CAS，存量系统零改造',
    '③ 国密合规：SM2签名/SM3哈希/SM4加密，满足等保三级要求',
    '④ 日志哈希链保护：SM3链式结构，任意篡改可检测',
    '⑤ 容器化部署：K8s编排，HPA自动扩缩容，22个Pod最小部署',
    '⑥ 国产化适配：麒麟OS + 达梦DM8 + 东方通TongWeb'
])

# ====== Slide 6: 实施计划 ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_text_box(slide, 0.8, 0.5, 11, 0.8, '实施计划（6个月）', 32, True, BLUE)
add_bullet_slide(slide, [
    '第一阶段（第1-2月）：基础平台搭建',
    '   - 用户管理+身份认证服务开发与部署',
    '   - 数据库/缓存/消息队列环境搭建',
    '第二阶段（第3-4月）：核心功能开发',
    '   - 权限管理+安全审计服务开发',
    '   - SSO协议集成（OAuth2/OIDC/SAML/CAS）',
    '第三阶段（第5-6月）：系统集成与上线',
    '   - 现有业务系统逐个对接',
    '   - 数据迁移+全量测试+试运行',
    '   - 管理员培训+用户培训'
])

# ====== Slide 7: 预算概览 ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_text_box(slide, 0.8, 0.5, 11, 0.8, '预算概览', 32, True, BLUE)
add_bullet_slide(slide, [
    '硬件/软件：230.5万元（服务器+数据库+中间件+容器平台）',
    '开发实施：38.5万元（18个功能模块，154人天）',
    '其他费用：75.9万元（项目管理+培训+质保运维+差旅）',
    '',
    '总投资：298.0万元（优惠报价）',
    '3年TCO：328.0万元（含第3年运维续保）'
])

# ====== Slide 8: 团队保障 ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_text_box(slide, 0.8, 0.5, 11, 0.8, '项目保障', 32, True, BLUE)
add_bullet_slide(slide, [
    '实施团队：项目经理1人 + 架构师1人 + 开发4人 + 测试2人 + 运维1人',
    '质保服务：验收后2年免费运维，7×12小时响应',
    '培训服务：管理员培训3天 + 用户培训2天 + 全套文档交付',
    '文档交付系统设计方案、接口规范、部署手册、运维手册、源代码'
])

# 保存
os.makedirs(os.path.dirname(OUTPUT), exist_ok=True)
prs.save(OUTPUT)
print(f"✅ PPT已生成: {OUTPUT}")
print(f"   共 {len(prs.slides)} 页幻灯片")
